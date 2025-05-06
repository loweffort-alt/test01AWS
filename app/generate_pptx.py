from pptx import Presentation
from pptx.util import Inches
import os


def crear_pptx(images, img_dir, name_pptx='triangles.pptx'):
    """
    Crea una presentación de PowerPoint con las imágenes proporcionadas.

    :param images: Lista de rutas a las imágenes.
    :param img_file: Ruta al archivo de imagen.
    :param name_pptx: Nombre del archivo de presentación a guardar.
    """
    # Crear una nueva presentación
    ppt = Presentation()

    # Agregar cada imagen a la presentación
    for image in images:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Layout en blanco

        # Ruta completa de la imagen
        img_path = os.path.join(img_dir, image)
        slide.shapes.add_picture(
            img_path, Inches(1), Inches(1), width=Inches(8))

    # Guardar la presentación
    ppt.save(name_pptx)
